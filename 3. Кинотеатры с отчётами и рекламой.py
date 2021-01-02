import datetime as dt

from copy import deepcopy
from collections import defaultdict

from docxtpl import DocxTemplate

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.dml.color import RGBColor


class BusyPlaceError(Exception):
    pass


class Cinema:
    def __init__(self, name):
        self.name = name
        self.halls = []

    def __str__(self):
        return f'Кинотеатр {self.name}: {len(self.halls)} залов'

    def __getitem__(self, item):
        return self.halls[item]

    def add_hall(self, hall):
        self.halls.append(hall)


class Hall:
    def __init__(self, cinema: Cinema):
        self.cinema = cinema
        self.plan = self.__make_halls_plan()
        self.num = len(self.cinema.halls) + 1
        self.movies = []

    def nothing(self):
        pass

    @staticmethod
    def __make_halls_plan():
        print('Задайте размер зала NxM, где N - количество рядов в зале,э')
        print('а M - количество мест в ряду. 0 < n <= 15, 0 < m <= 30')

        n = m = 0
        while not (0 < n <= 15 and 0 < m <= 30):
            try:
                n, m = map(int, input('Введите размер зала: ').split())
            except ValueError:
                print_error('Введены некорректные размеры')
                continue

            if not (0 < n <= 15 and 0 < m <= 30):
                print_error('Введены недопустимые размеры')
        else:
            return [[str(i).rjust(2) for i in range(1, m + 1)] for _ in range(n)]

    @staticmethod
    def choose_movie(movies):
        if len(movies) == 1:
            return movies[0]

        print(f'В данном зале проходит несколько фильмов с названием {movies[0].name}')
        print(*[f'{i}. {film.name} {film.start}-{film.end}' for i, film in enumerate(movies, 1)],
              sep='\n')

        while True:
            try:
                user_ans = input('Выберите номер сеанса, на который хотите купить билет или '
                                 'посмотреть места\n')
                assert 0 < int(user_ans) <= len(movies)
                return movies[int(user_ans) - 1]
            except (ValueError, IndexError, AssertionError):
                print_error('Введите корректный номер сеанса')

    def print_plan(self):
        print(' ' * 3 + '-' * (len(self.plan[0]) - 5) * 4)
        for i, row in enumerate(self.plan, 1):
            print(str(i).ljust(2) + '-' + ' '.join(row))

    def add_movie(self, movie):
        for film in self.movies:
            if film == movie:
                print_warning(f'В это время уже идет фильм {film.name} с {film.start} до {film.end}')
                return 'Error'

        self.movies.append(movie)
        print(f'Добавляю фильм {movie.name} с {movie.start} до {movie.end}')

    def get_movie(self, movie):
        suitable_movies = []
        for film in self.movies:
            if film.name == movie:
                suitable_movies.append(film)

        return self.choose_movie(suitable_movies) if suitable_movies else None


class Movie:
    def __init__(self, hall: Hall):
        self.name = input('Введите название фильма\n')

        self.hall = hall
        self.halls_plan = deepcopy(hall.plan)
        self.num = len(self.hall.movies) + 1

        self.start = self.end = None
        self.set_time()

        self.orders = []

    def __eq__(self, other):
        return (other.start <= self.start < other.end or other.start < self.end < other.end or
                (self.start < other.start and self.end > other.end) or
                (self.start > other.start and self.end < other.end))

    def set_time(self):
        movie_time = []

        for i, time in enumerate((self.start, self.end)):
            while time is None:
                try:
                    time = dt.time(*map(int, input(f'Введите время {"начала" if i == 0 else "конца"}'
                                                   f' кино в формате "hh mm"\n').split()))
                    movie_time.append(time)
                except (ValueError, TypeError):
                    print_error('Введена некорректное время фильма. Формат ввода "hh mm".')

        self.start, self.end = sorted(movie_time)

    def show_hall(self):
        print(' ' * 3 + '-' * (len(self.halls_plan[0]) - 5) * 4)
        for i, row in enumerate(self.halls_plan, 1):
            print(str(i).ljust(2) + '-' + ' '.join(row))

    def buy_tickets(self):
        number = input('Выберите количество билетов (от 1 до 5)\n')
        while number not in map(str, range(1, 6)):
            number = input('Выберите количество билетов (от 1 до 5)\n')
        number = int(number)

        tickets = []
        for _ in range(number):
            print('Выберите места (точкой помечаются забронированные места)')
            print('Для выбора места введи ряд и номер места через пробел (например 1 2)')
            self.show_hall()

            ticket = None
            while ticket is None:
                try:
                    n, m = map(int, input().split())
                    assert n > 0 and m > 0

                    if self.halls_plan[n - 1][m - 1] == ' .':
                        raise BusyPlaceError

                    self.halls_plan[n - 1][m - 1] = ' .'
                    ticket = (n, m)
                    tickets.append(ticket)
                    print(f'Бронирую место {n} {m}')
                except (ValueError, AssertionError, IndexError):
                    print_error('Выберите корректное место')
                except BusyPlaceError:
                    print_warning('Данное место уже занято')
        self.show_hall()

        name = input('Введите на чье имя записать билеты\n')
        self.orders.append((name, tickets, len(tickets)))


COMMANDS = ('add_cinema', 'add_hall', 'add_movie', 'check_movie', 'show_hall', 'buy_ticket', 'exit',
            'help', 'check_orders', 'c', 'h', 'm', 'cm', 'sh', 'bt', 'e', 'co')

cinemas = {}
all_orders = defaultdict(list)

reports_amount = 0
presentation_amount = 0


def interface():
    print('Список возможных комманд:')
    print()
    print('Интерфейс для владельца:')
    print('add_cinema/c <name>                     -- добавить кинотеатр с именем name.')
    print('add_hall/h <cinema>                     -- добавить зал, в кинотеатр cinema. Залы будут')
    print('                                           автоматически нумераться, начиная с 1.')
    print('add_movie/m <cinema> <number>           -- добавить фильм в расписание зала number в')
    print('                                           кинотеатре cinema (создать рекламный буклет).')
    print('check_orders/co                         -- Показать все заказы по всем кинотеатрам')
    print('                                           (сгенерировать отчет по посещаемости).')
    print()
    print('Интерфей для пользователя:')
    print('check_movie/cm <name>                   -- вывести список всех кинотеатров и залов,')
    print('                                           в которых есть свободные места на кино name.')
    print('show_hall/sh <cinema> <number> <movie>  -- вывести места зала в кинотеатре cinema')
    print('                                           в зале number на фильм movie.')
    print('buy_ticket/bt <cinema> <number> <movie> -- купить билет(ы) на сеанс.')
    print()
    print('exit/e                                  -- выход.')
    print('help            -- вывести подсказку (через / обозначается сокращенный вариант команды).')
    print('Примечание: 1. символы "<>" используются для обозначения аргумента, вводить данные')
    print('               следует без них (например add_cinema Киргизия_Новая).')
    print('            2. Названия кинотеатров и фильмов следует вводить в одно слово (используйте')
    print('               "_" для разделения слов в названиях).')


def main():
    command = input('Введите команду\n')
    while not command or command.split()[:1][0] not in COMMANDS:
        print_error('Введена некорректная команда')
        command = input('Введите команду\n')
    command = command.split()

    if command[0] in ('add_cinema', 'c') and len(command) == 2:
        name = command[1]

        if cinemas.get(name) is not None:
            print_warning('Кинотеатр с таким именем уже существует')
            return True

        print(f'Добавлен кинотеатр {name}')
        cinemas[name] = Cinema(name)
    elif command[0] in ('add_hall', 'h') and len(command) == 2:
        cinema = check_cinema(command[1])
        if cinema is None:
            return True

        hall = Hall(cinema)
        print(f'Добавляю зал в кинотеатр {cinema.name}')
        cinema.add_hall(hall)
    elif command[0] in ('add_movie', 'm') and len(command) == 3:
        cinema = check_cinema(command[1])
        hall = check_hall(cinema, command[2])
        if cinema is None or hall is None:
            return True

        movie = Movie(cinema[hall])
        req = cinema[hall].add_movie(movie)
        if req == 'Error':
            return True

        user_ans = input('Сгенерировать рекламный буклет для фильма? (да, нет)\n')
        while user_ans not in ('да', 'нет'):
            print_error('Введите корректный ответ (да или нет)')
            user_ans = input()

        if user_ans == 'да':
            start_time = f"{movie.start.hour:02d}:{movie.start.minute:02d}"
            end_time = f"{movie.end.hour:02d}:{movie.end.minute:02d}"

            req = generate_adv_pres(cinema.name, hall, movie.name, start_time, end_time)
            while req == 'Error':
                req = generate_adv_pres(cinema.name, hall, movie.name, start_time, end_time)

            print('Рекламный буклет создан')
    elif command[0] in ('check_orders', 'co') and len(command) == 1:
        print_all_orders()

        user_ans = input('Сгенерировать отчет о посещаемости кинотеатров? (да, нет)\n')
        while user_ans not in ('да', 'нет'):
            print_error('Введите корректный ответ (да или нет)')
            user_ans = input()

        if user_ans == 'да':
            generate_docx_report()
            print('Отчет создан')
    elif command[0] in ('check_movie', 'cm') and len(command) == 2:
        print('Хотите искать сеансы, на которых будет определенное количество мест рядом?')
        print('Введите 0 если не хотите, иначе введите количество мест рядом (от 1 до 5)')

        number = input()
        while number not in map(str, range(6)):
            print_error('Введите корректное число мест')
            number = input()

        check_movie(command[1], int(number))
    elif command[0] in ('show_hall', 'sh') and len(command) == 4:
        cinema = check_cinema(command[1])
        hall = check_hall(cinema, command[2])
        if cinema is None or hall is None:
            return True

        movie = cinema[hall].get_movie(command[3])
        if movie is None:
            print_warning('Данного фильма нет в прокате в данном зале')
            return True

        print(f'Вывожу места зала номер {hall + 1} в кинотеатре {cinema} на фильм {movie.name}')
        movie.show_hall()
    elif command[0] in ('buy_ticket', 'bt') and len(command) == 4:
        cinema = check_cinema(command[1])
        hall = check_hall(cinema, command[2])
        if cinema is None or hall is None:
            return True

        movie = cinema[hall].get_movie(command[3])
        if movie is None:
            print_warning('Данного фильма нет в прокате в данном зале')
            return True
        movie.buy_tickets()
    elif command[0] in ('exit', 'e') and len(command) == 1:
        return False
    elif command[0] == 'help' and len(command) == 1:
        interface()
    else:
        print_error('Введена некорректная команда')

    return True


def check_cinema(cinema):
    if cinemas.get(cinema) is None:
        print_warning('Данного кинотеатра не существует')

        print('Имеющиеся кинотеатры:')
        for cinema in cinemas:
            print(cinemas[cinema])

        return None
    return cinemas[cinema]


def check_hall(cinema, hall_number):
    if cinema is None:
        return
    try:
        assert int(hall_number) > 0
        cinema[int(hall_number) - 1].nothing()
        return int(hall_number) - 1
    except (ValueError, IndexError, AssertionError):
        print_warning(f'Данного зала в кинотеатре {cinema.name} не сущетсвует')
        return None


def print_all_orders():
    for cinema in cinemas:
        print(f'Кинотеатр {cinema}')
        for i, hall in enumerate(cinemas[cinema].halls, 1):
            print(f'\tЗал номер {i}')
            for j, movie in enumerate(hall.movies, 1):
                print(f'\t\t{j}. {movie.name}: {movie.start} - {movie.end}')

                if not movie.orders:
                    print('\t\t\tНа данный сеанс билеты еще не бронировались')
                for h, (name, tickets, amount) in enumerate(movie.orders, 1):
                    print(f'\t\t\t{j}.{h} {name} - {amount} билетов забронировано: '
                          f'{", ".join([str(el) for el in tickets])}')


def generate_docx_report():
    global reports_amount

    doc = DocxTemplate("tpl.docx")
    context = {'cinemas': cinemas}
    doc.render(context)
    doc.save(f"report{reports_amount if reports_amount != 0 else ''}.docx")

    reports_amount += 1


def generate_adv_pres(cinema_name, hall_num, movie_name, movie_start, movie_end):
    global presentation_amount

    def place_text(txt: str, pos: tuple, color=(0, 0, 0), pt_size=24,
                   bold=False, italic=False, alignment=PP_ALIGN.LEFT):
        text_box = slide.shapes.add_textbox(*pos)
        text_frame = text_box.text_frame

        p = text_frame.add_paragraph()
        run = p.add_run()
        font = run.font

        font.size = Pt(pt_size)
        font.bold = bold
        font.italic = italic
        p.alignment = alignment
        font.color.rgb = RGBColor(*color)
        run.text = txt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.shapes.add_picture('new_film.jpg', 0, 0, prs.slide_width, prs.slide_height)

    text = f"Смотрите {movie_name} в {hall_num + 1} зале кинотеатра {cinema_name}\n" \
           f"Время сеанса: {movie_start} - {movie_end}"
    place_text("В прокат вышел новый фильм!", (0, 0, prs.slide_width, 566_738),
               (255, 255, 0), 44, True, alignment=PP_ALIGN.CENTER)
    place_text(text, (0, 4_800_600, prs.slide_width, 566_738), (255, 255, 9), 30,
               alignment=PP_ALIGN.CENTER)

    name = f"adv_booklet{presentation_amount if presentation_amount != 0 else ''}.pptx"
    try:
        prs.save(name)
    except PermissionError:
        text = f"Файл {name} невозможно редактировать, тк он открыт в настоящий момент.\n" \
               f"Пожалуйста закройте файл и введите 'ок' для продолжения работы"
        print_error(text)
        user_ans = input()
        while user_ans != 'ок':
            print_error("ОШИБКА! " + text)
            user_ans = input()

        return 'Error'

    presentation_amount += 1


def check_movie(name, number):
    number = number or 1
    suitable_movies = []
    is_suitable_movie = False

    for cinema in cinemas:
        for i, hall in enumerate(cinemas[cinema].halls, 1):
            for j, movie in enumerate(hall.movies, 1):
                if movie.name != name:
                    continue
                is_suitable_movie = True

                if there_are_enough_places(movie.halls_plan, number):
                    suitable_movies.append((cinemas[cinema], hall, movie))

    if not is_suitable_movie:
        print_warning('Такого фильма нигде нет в прокате')
        return

    if not suitable_movies:
        print_warning('Нет подходящих сеансов')
    else:
        print('Найдены подходящие сеансы:')
        for i, (cinema, hall, movie) in enumerate(suitable_movies, 1):
            print(f'{i}. Кинотеатр {cinema.name}, зал номер {cinema.halls.index(hall) + 1}')
            print(f'{len(str(i)) * " "}  фильм {movie.name}, сеанс {movie.start}-{movie.end}')


def there_are_enough_places(hall_plan, number):
    for row in hall_plan:
        places = ''.join([el if el == ' .' else '+-' for el in row])
        if '+-' * number in places:
            return True

    return False


def print_error(text):
    print("\033[31m{}\033[0m" .format(text))


def print_warning(text):
    print("\033[33m{}\033[0m".format(text))


if __name__ == '__main__':
    print('Вас приветствует Билетная система сети кинотеатров Яндекс.Кино.')
    interface()
    while main():
        pass
    print('Спасибо за использование нашей сети Яндекс.Кино')

"""
c a
h a
4 4
h a
6 6

m a 1
asdf
12 00
13 00
да

m a 2
asdf
13 00
14 00
да

bt a 1 asdf
4
1 1
1 2
1 3
1 4
Олеся

bt a 2 asdf
3
2 1
2 2
2 3
Никита Жуков

m a 1
Шаман_кинг
13 00
14 00
нет

bt a 1 Шаман_кинг
2
2 3
3 2
Жуков Сергей

c Шангал
h Шангал
10 10

m Шангал 1
Шаман_кинг
11 00
12 00
да

"""
